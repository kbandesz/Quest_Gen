import io, math, re
from typing import Tuple, List, Union
from pypdf import PdfReader
import mammoth  # for docx to text
from pptx import Presentation  # for pptx to text (python-pptx)


# PDF
def _read_pdf(f) -> str:
    reader = PdfReader(f)
    texts=[]
    for p in reader.pages:
        try:
            texts.append(p.extract_text() or "")
        except Exception:
            continue
    return "\n".join(texts)

# DOCX
def _read_docx(f) -> str:
    b = f.read() if hasattr(f,"read") else f
    result = mammoth.extract_raw_text(io.BytesIO(b))
    return result.value or ""

# TXT
def _read_txt(f)->str:
    return f.read().decode("utf-8", errors="ignore") if hasattr(f,"read") else f.decode("utf-8", errors="ignore")

# PPTX
def _read_pptx(f) -> str:
    """
    Extract text from all slides (shapes with text frames) and speaker notes.
    """
    b = f.read() if hasattr(f, "read") else f
    prs = Presentation(io.BytesIO(b))
    chunks = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_bits = []

        # Slide body text
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                # join all paragraphs/runs into a single string per shape
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                if text:
                    slide_bits.append(text)

        # Speaker notes (if any)
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = "\n".join(
                p.text for p in slide.notes_slide.notes_text_frame.paragraphs if p.text
            )
            if notes_text:
                slide_bits.append(f"[Notes]\n{notes_text}")

        if slide_bits:
            chunks.append(f"[Slide {idx}]\n" + "\n".join(slide_bits))

    return "\n\n".join(chunks)


# Consistent newline char, reduce multiple whitespace chars, remove trailing spaces
def _normalize(text: str) -> str:
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()

# Extract text from a single file
def _extract_single(uploaded_file) -> str:
    try:
        name = uploaded_file.name.lower()
        if name.endswith(".pdf"):
            return _normalize(_read_pdf(uploaded_file))
        if name.endswith(".docx"):
            return _normalize(_read_docx(uploaded_file))
        if name.endswith(".pptx"):
            return _normalize(_read_pptx(uploaded_file))
        if name.endswith(".txt"):
            return _normalize(_read_txt(uploaded_file))
        raise ValueError("Unsupported file type")
    except Exception as e:
        raise ValueError(f"Failed to parse the file: {uploaded_file.name}. Please check the file type and ensure it is not password-protected or corrupted.")

def extract_text_and_tokens(uploaded_files: Union[List, object]) -> Tuple[str, int]:
    """
    Accepts a single uploaded file or a list of files.
    Returns a single combined text and a token estimate.
    """
    if not uploaded_files:
        return "", 0

    # Support both single-file and multi-file callers
    files = uploaded_files if isinstance(uploaded_files, list) else [uploaded_files]

    # Combine in the order provided, with a clear separator between files
    parts: List[str] = []
    for f in files:
        parts.append(_extract_single(f))
    combined = ("\n\n----- FILE BREAK -----\n\n").join(p for p in parts if p)

    # heuristic token estimate: 1 token ≈ 0.75 words  -> tokens ≈ words / 0.75
    words = len(re.findall(r"\S+", combined))
    tokens_est = int(math.ceil(words / 0.75))
    return combined, tokens_est
